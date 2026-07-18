"""Bounded extraction for the explicit set of AI-indexable file formats."""

from __future__ import annotations

import io
import multiprocessing
import os
from pathlib import PurePosixPath
import zipfile
from dataclasses import dataclass


_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif", ".webp")
_TEXT_EXTS = (
    ".txt", ".md", ".csv", ".tsv", ".json", ".jsonl", ".html", ".htm",
    ".xml", ".yaml", ".yml", ".log",
)
_SUPPORTED_EXTS = (".pdf", ".docx", ".xlsx", ".xlsm", ".pptx", ".rtf", *_IMAGE_EXTS, *_TEXT_EXTS)


class UnsupportedDocumentError(ValueError):
    """The original can be stored, but is unsupported or ambiguous for AI."""


@dataclass(frozen=True)
class ExtractionLimits:
    max_source_bytes: int = 50 * 1024 * 1024
    max_extracted_characters: int = 5_000_000
    max_pdf_pages: int = 500
    max_ocr_pages: int = 100
    max_sheets: int = 200
    max_rows_per_sheet: int = 200_000
    max_cells: int = 2_000_000
    max_slides: int = 2_000
    max_shapes: int = 100_000
    max_image_pixels: int = 80_000_000
    ocr_timeout_seconds: int = 60
    max_archive_entries: int = 20_000
    max_archive_uncompressed_bytes: int = 250 * 1024 * 1024
    max_archive_member_bytes: int = 100 * 1024 * 1024
    max_archive_compression_ratio: int = 200
    total_timeout_seconds: int = 180
    max_worker_memory_bytes: int = 1024 * 1024 * 1024


DEFAULT_EXTRACTION_LIMITS = ExtractionLimits()


def supports_filename(filename: str) -> bool:
    return (filename or "").lower().endswith(_SUPPORTED_EXTS)


def extract_text_isolated(
    filename: str, data: bytes, limits: ExtractionLimits | None = None,
) -> str:
    """Extract in a killable worker with wall-clock and OS resource bounds."""

    limits = limits or DEFAULT_EXTRACTION_LIMITS
    if len(data) > limits.max_source_bytes:
        raise ValueError("Document exceeds the extraction source-size limit.")
    context = multiprocessing.get_context("spawn")
    receiver, sender = context.Pipe(duplex=False)
    process = context.Process(
        target=_isolated_worker,
        args=(sender, filename, data, limits),
        name="onebrain-document-extractor",
        daemon=True,
    )
    process.start()
    sender.close()
    try:
        if not receiver.poll(max(1, int(limits.total_timeout_seconds))):
            process.terminate()
            process.join(timeout=5)
            if process.is_alive() and hasattr(process, "kill"):
                process.kill()
                process.join(timeout=5)
            raise TimeoutError("Document extraction exceeded the configured time limit.")
        try:
            status, value = receiver.recv()
        except EOFError as exc:
            raise RuntimeError("Document extractor exited without a result.") from exc
    finally:
        receiver.close()
        process.join(timeout=5)
        if process.is_alive():
            process.terminate()
            process.join(timeout=5)
    if status == "ok":
        return value
    if status == "unsupported":
        raise UnsupportedDocumentError(value)
    if status == "value":
        raise ValueError(value)
    raise RuntimeError(value)


def _isolated_worker(sender, filename: str, data: bytes, limits: ExtractionLimits) -> None:
    try:
        if os.name != "nt":
            try:
                import resource

                memory = max(256 * 1024 * 1024, int(limits.max_worker_memory_bytes))
                resource.setrlimit(resource.RLIMIT_AS, (memory, memory))
                cpu = max(1, int(limits.total_timeout_seconds))
                resource.setrlimit(resource.RLIMIT_CPU, (cpu, cpu + 1))
            except (ImportError, OSError, ValueError):
                # The parent wall-clock kill remains mandatory on platforms that
                # cannot install an address-space/CPU rlimit.
                pass
        sender.send(("ok", extract_text(filename, data, limits)))
    except UnsupportedDocumentError as exc:
        sender.send(("unsupported", str(exc)[:1000]))
    except ValueError as exc:
        sender.send(("value", str(exc)[:1000]))
    except BaseException as exc:
        sender.send(("runtime", f"Document extraction failed: {exc}"[:1000]))
    finally:
        sender.close()


def extract_text(filename: str, data: bytes, limits: ExtractionLimits | None = None) -> str:
    limits = limits or DEFAULT_EXTRACTION_LIMITS
    if len(data) > limits.max_source_bytes:
        raise ValueError("Document exceeds the extraction source-size limit.")
    name = (filename or "").lower()
    if not supports_filename(name):
        raise UnsupportedDocumentError("This file type is stored but unsupported for AI indexing.")
    if name.endswith(".pdf"):
        if not data.startswith(b"%PDF"):
            raise UnsupportedDocumentError("The file extension does not match PDF content.")
        return _bounded(_pdf(data, limits), limits)
    if name.endswith(".docx"):
        _require_safe_zip(data, limits)
        return _bounded(_docx(data), limits)
    if name.endswith((".xlsx", ".xlsm")):
        _require_safe_zip(data, limits)
        return _bounded(_xlsx(data, limits), limits)
    if name.endswith(".pptx"):
        _require_safe_zip(data, limits)
        return _bounded(_pptx(data, limits), limits)
    if name.endswith(".rtf"):
        if not data.lstrip().startswith(b"{\\rtf"):
            raise UnsupportedDocumentError("The file extension does not match RTF content.")
        return _bounded(_rtf(data), limits)
    if name.endswith(_IMAGE_EXTS):
        return _bounded(_ocr_image(data, limits), limits)
    if b"\x00" in data[:8192]:
        raise UnsupportedDocumentError("Binary content cannot be indexed as plain text.")
    return _bounded(_plain_text(data), limits)


def _plain_text(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        # Known text extensions may use legacy Western encodings. The extension
        # allowlist and NUL guard above prevent arbitrary binaries reaching this.
        return data.decode("latin-1")


def _pdf(data: bytes, limits: ExtractionLimits) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("PDF support needs pypdf") from exc
    reader = PdfReader(io.BytesIO(data))
    if len(reader.pages) > limits.max_pdf_pages:
        raise ValueError("PDF exceeds the configured page limit.")
    text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
    if len(text.strip()) >= 20:
        return text
    ocr = _ocr_pdf(data, limits)
    return ocr if ocr.strip() else text


def _docx(data: bytes) -> str:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("Word (.docx) support needs python-docx") from exc
    document = docx.Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx(data: bytes, limits: ExtractionLimits) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Excel (.xlsx) support needs openpyxl") from exc
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    if len(workbook.worksheets) > limits.max_sheets:
        workbook.close()
        raise ValueError("Spreadsheet exceeds the configured sheet limit.")
    parts: list[str] = []
    cell_count = 0
    for sheet in workbook.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), 1):
            if row_index > limits.max_rows_per_sheet:
                workbook.close()
                raise ValueError("Spreadsheet exceeds the configured row limit.")
            cells = [str(cell) for cell in row if cell is not None]
            cell_count += len(cells)
            if cell_count > limits.max_cells:
                workbook.close()
                raise ValueError("Spreadsheet exceeds the configured cell limit.")
            if cells:
                parts.append(" | ".join(cells))
    workbook.close()
    return "\n".join(parts)


def _pptx(data: bytes, limits: ExtractionLimits) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PowerPoint (.pptx) support needs python-pptx") from exc
    presentation = Presentation(io.BytesIO(data))
    if len(presentation.slides) > limits.max_slides:
        raise ValueError("Presentation exceeds the configured slide limit.")
    parts: list[str] = []
    shape_count = 0
    for index, slide in enumerate(presentation.slides, 1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            shape_count += 1
            if shape_count > limits.max_shapes:
                raise ValueError("Presentation exceeds the configured shape limit.")
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def _rtf(data: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as exc:
        raise RuntimeError("RTF support needs striprtf") from exc
    return rtf_to_text(_plain_text(data))


def _ocr_image(data: bytes, limits: ExtractionLimits) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError("Image OCR needs pytesseract + Pillow") from exc
    try:
        image = Image.open(io.BytesIO(data))
        if image.width * image.height > limits.max_image_pixels:
            raise ValueError("Image exceeds the configured pixel limit.")
        return pytesseract.image_to_string(image, timeout=limits.ocr_timeout_seconds)
    except ValueError:
        raise
    except Exception as exc:
        raise RuntimeError(f"Image OCR failed (is Tesseract installed?): {exc}") from exc


def _ocr_pdf(data: bytes, limits: ExtractionLimits) -> str:
    try:
        import fitz
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""
    try:
        output = []
        with fitz.open(stream=data, filetype="pdf") as document:
            for page_index, page in enumerate(document):
                if page_index >= limits.max_ocr_pages:
                    raise ValueError("Scanned PDF exceeds the configured OCR page limit.")
                pixmap = page.get_pixmap(dpi=200)
                image = Image.open(io.BytesIO(pixmap.tobytes("png")))
                if image.width * image.height > limits.max_image_pixels:
                    raise ValueError("Rendered PDF page exceeds the configured pixel limit.")
                output.append(pytesseract.image_to_string(image, timeout=limits.ocr_timeout_seconds))
        return "\n\n".join(output)
    except ValueError:
        raise
    except Exception:
        return ""


def _bounded(text: str, limits: ExtractionLimits) -> str:
    if len(text) > limits.max_extracted_characters:
        raise ValueError("Document exceeds the configured extracted-text limit.")
    return text


def _require_safe_zip(data: bytes, limits: ExtractionLimits) -> None:
    if not data.startswith(b"PK\x03\x04"):
        raise UnsupportedDocumentError("The Office file extension does not match ZIP content.")
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = archive.infolist()
    except (zipfile.BadZipFile, OSError) as exc:
        raise UnsupportedDocumentError("The Office archive is malformed.") from exc
    if len(entries) > limits.max_archive_entries:
        raise ValueError("Office archive exceeds the configured entry limit.")
    expanded = 0
    for entry in entries:
        normalized = entry.filename.replace("\\", "/")
        path = PurePosixPath(normalized)
        if path.is_absolute() or ".." in path.parts:
            raise UnsupportedDocumentError("Office archive contains an unsafe member path.")
        if entry.flag_bits & 0x1:
            raise UnsupportedDocumentError("Encrypted Office archives cannot be indexed.")
        if entry.file_size > limits.max_archive_member_bytes:
            raise ValueError("Office archive member exceeds the configured expansion limit.")
        expanded += entry.file_size
        if expanded > limits.max_archive_uncompressed_bytes:
            raise ValueError("Office archive exceeds the configured expansion limit.")
        if entry.file_size:
            if entry.compress_size <= 0:
                raise ValueError("Office archive has an unsafe compression ratio.")
            if entry.file_size / entry.compress_size > limits.max_archive_compression_ratio:
                raise ValueError("Office archive has an unsafe compression ratio.")
